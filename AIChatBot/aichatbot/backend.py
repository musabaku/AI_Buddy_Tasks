import os
from flask import Flask, render_template, request, redirect, url_for, flash, session
from flask_session import Session  # New import
import PyPDF2
from pptx import Presentation
import google.generativeai as genai
from werkzeug.utils import secure_filename

genai.configure(api_key=os.getenv("GEMINI_API_KEY", "AIzaSyDRmoe4y30uppKWDVgCzjVHejKzhj8wxYI"))

app = Flask(__name__, template_folder='frontend')
app.secret_key = "some_secret_key"

app.config["SESSION_TYPE"] = "filesystem"
Session(app)  # Initialize the session

UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def set_document_text(text):
    session["document_text"] = text

def get_document_text():
    return session.get("document_text", "")

def generate_text(prompt):
    try:
        model = genai.GenerativeModel("models/gemini-1.5-flash")
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"Error calling Google Generative AI: {e}"

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['GET', 'POST'])
def upload():
    if request.method == 'POST':
        if 'file' not in request.files:
            flash("No file part in the request.")
            return redirect(request.url)
        file = request.files['file']
        if file.filename == '':
            flash("No file selected.")
            return redirect(request.url)
        if file:
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            # Process file based on extension.
            if filename.lower().endswith(".pdf"):
                try:
                    with open(filepath, "rb") as f:
                        reader = PyPDF2.PdfReader(f)
                        text = ""
                        for page in reader.pages:
                            text += page.extract_text() or ""
                    set_document_text(text)
                except Exception as e:
                    flash(f"Error processing PDF: {e}")
                    return redirect(request.url)
            elif filename.lower().endswith((".ppt", ".pptx")):
                try:
                    prs = Presentation(filepath)
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    set_document_text(text)
                except Exception as e:
                    flash(f"Error processing PowerPoint file: {e}")
                    return redirect(request.url)
            else:
                flash("Unsupported file type. Please upload a PDF or PPT/PPTX file.")
                return redirect(request.url)

            flash("File uploaded and processed successfully!")
            return redirect(url_for("document_query"))
    return render_template('upload.html')

@app.route('/document_query', methods=['GET', 'POST'])
def document_query():
    ai_response = None
    if request.method == 'POST':
        user_input = request.form.get('user_input')
        focus = request.form.get('focus')
        document_text = get_document_text()
        if not document_text:
            flash("No document content available. Please upload a file first.")
            return redirect(url_for('upload'))

        if focus:
            prompt = (
                f"Based on the following document content:\n\n{document_text}\n\n"
                f"Generate a {focus.lower()} answering the following query:\n{user_input}"
            )
        else:
            prompt = (
                f"Based on the following document content:\n\n{document_text}\n\n"
                f"Answer the following query:\n{user_input}"
            )
        ai_response = generate_text(prompt)
    return render_template('document_query.html', response=ai_response)

@app.route('/chat', methods=['GET', 'POST'])
def chat():
    ai_response = None
    if request.method == 'POST':
        user_input = request.form.get('user_input')
        document_text = get_document_text()
        if document_text:
            prompt = f"Based on the following document content:\n\n{document_text}\n\nUser query: {user_input}"
        else:
            prompt = user_input
        ai_response = generate_text(prompt)
    return render_template('chat.html', response=ai_response)

@app.route('/generate', methods=['GET', 'POST'])
def generate():
    ai_response = None
    if request.method == 'POST':
        generation_type = request.form.get('generation_type')
        additional_info = request.form.get('additional_info')
        if generation_type == "lesson_plan":
            prompt = (
                f"Create a detailed lesson plan for a subject. "
                f"Include objectives, activities, and assessments. {additional_info}"
            )
        elif generation_type == "timetable":
            prompt = (
                f"Generate an engaging timetable for a school day. "
                f"Include class times, breaks, and extra-curricular activities. {additional_info}"
            )
        elif generation_type == "quiz":
            prompt = (
                f"Design a quiz with diverse question types (multiple choice, short answers, true/false) "
                f"on a given topic. {additional_info}"
            )
        elif generation_type == "educational_material":
            prompt = (
                f"Develop creative educational materials that can be used in a classroom setting. "
                f"Include interactive activities and clear explanations. {additional_info}"
            )
        else:
            prompt = additional_info

        ai_response = generate_text(prompt)
    return render_template('generate.html', response=ai_response)

if __name__ == '__main__':
    app.run(debug=True)
