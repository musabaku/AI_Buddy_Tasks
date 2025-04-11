import os
from flask import Flask, render_template, request, redirect, url_for, flash
import PyPDF2
from pptx import Presentation
import google.generativeai as genai  # Updated import

# Configure the Generative AI library with your API key.
# (Replace "GEMINI_API_KEY" with your environment variable name as needed.)
genai.configure(api_key=os.getenv("GEMINI_API_KEY", "AIzaSyDRmoe4y30uppKWDVgCzjVHejKzhj8wxYI"))

app = Flask(__name__, template_folder='frontend')
app.secret_key = "some_secret_key"

UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Global variable to hold extracted document text (for demo purposes)
document_text = ""

# Function to generate text using the Google Generative AI library.
# def generate_text(prompt):
#     try:
#         # Call the generate_text function with your model and prompt.
#         # Note: Depending on the library version, the function may return a dict.
#         response = genai.generate_content(
#             model="models/gemini-1.5-flash",  # or an available model name if different
#             prompt=prompt,
#             temperature=0.7,         # adjust the temperature as needed
#         )
#         # Assuming the response is structured as:
#         # {
#         #   "candidates": [
#         #       {"output": "Generated text here."}
#         #   ]
#         # }
#         if "candidates" in response and len(response["candidates"]) > 0:
#             return response["candidates"][0].get("output", "No output returned.")
#         else:
#             return "No candidates returned in response."
#     except Exception as e:
#         return f"Error calling Google Generative AI: {e}"
def generate_text(prompt):
    try:
        model = genai.GenerativeModel("models/gemini-1.5-flash")
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"Error calling Google Generative AI: {e}"
# Home page route
@app.route('/')
def index():
    return render_template('index.html')

# Upload route: GET shows the form; POST processes the uploaded file
@app.route('/upload', methods=['GET', 'POST'])
def upload():
    global document_text
    if request.method == 'POST':
        if 'file' not in request.files:
            flash("No file part in the request.")
            return redirect(request.url)
        file = request.files['file']
        if file.filename == '':
            flash("No file selected.")
            return redirect(request.url)
        if file:
            filename = file.filename
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)

            # Process file based on extension
            if filename.lower().endswith(".pdf"):
                try:
                    with open(filepath, "rb") as f:
                        reader = PyPDF2.PdfReader(f)
                        text = ""
                        for page in reader.pages:
                            text += page.extract_text() or ""
                        document_text = text
                except Exception as e:
                    flash(f"Error processing PDF: {e}")
                    return redirect(request.url)
            elif filename.lower().endswith((".ppt", ".pptx")):
                try:
                    prs = Presentation(filepath)
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    document_text = text
                except Exception as e:
                    flash(f"Error processing PowerPoint file: {e}")
                    return redirect(request.url)
            else:
                flash("Unsupported file type. Please upload a PDF or PPT/PPTX file.")
                return redirect(request.url)

            flash("File uploaded and processed successfully!")
            return redirect(url_for("chat"))
    return render_template('upload.html')

# Chat route: GET shows chat interface; POST handles queries
@app.route('/chat', methods=['GET', 'POST'])
def chat():
    ai_response = None
    if request.method == 'POST':
        user_input = request.form.get('user_input')
        # Build the prompt using the extracted document text (if available)
        prompt = user_input
        if document_text:
            prompt = f"Based on the following document content:\n\n{document_text}\n\nUser query: {user_input}"
        ai_response = generate_text(prompt)
    return render_template('chat.html', response=ai_response)

if __name__ == '__main__':
    # Optionally, you can set the API key via code if not using environment variables.
    # os.environ['GEMINI_API_KEY'] = "AIzaSyDRmoe4y30uppKWDVgCzjVHejKzhj8wxYI"
    app.run(debug=True)
